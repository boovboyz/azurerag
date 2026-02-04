from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import openpyxl

def extract_text(path):
    if path.endswith(".pdf"):
        return "\n".join(p.extract_text() or "" for p in PdfReader(path).pages)

    if path.endswith(".docx"):
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs)

    if path.endswith(".pptx"):
        prs = Presentation(path)
        return "\n".join(
            shape.text for slide in prs.slides
            for shape in slide.shapes if hasattr(shape, "text")
        )

    if path.endswith(".xlsx"):
        wb = openpyxl.load_workbook(path)
        text = ""
        for sheet in wb:
            for row in sheet.iter_rows(values_only=True):
                text += " ".join(str(c) for c in row if c) + "\n"
        return text

    return ""
